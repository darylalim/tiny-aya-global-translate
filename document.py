from __future__ import annotations

import io
from typing import Any

# -- Helpers -------------------------------------------------------------------


def _replace_paragraph_text(para: Any, text: str) -> None:
    """Replace all text in a paragraph, preserving the first run's formatting."""
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    elif text.strip():
        para.add_run().text = text


# -- DOCX ---------------------------------------------------------------------


def _iter_docx_paragraphs(doc: Any) -> list[Any]:
    """Return all paragraphs in document body and tables, deduped."""
    seen: set[int] = set()
    paragraphs: list[Any] = []
    for para in doc.paragraphs:
        pid = id(para._element)
        if pid not in seen:
            seen.add(pid)
            paragraphs.append(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    pid = id(para._element)
                    if pid not in seen:
                        seen.add(pid)
                        paragraphs.append(para)
    return paragraphs


def extract_segments_docx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    return [para.text for para in _iter_docx_paragraphs(doc)]


def rebuild_document_docx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild a DOCX file with translated text replacing original segments."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    for i, para in enumerate(_iter_docx_paragraphs(doc)):
        if i >= len(translations):
            break
        _replace_paragraph_text(para, translations[i])
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# -- PPTX ---------------------------------------------------------------------


def _iter_pptx_paragraphs(prs: Any) -> list[Any]:
    """Return all paragraphs across all slides, shapes, and tables."""
    paragraphs: list[Any] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    paragraphs.append(para)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            paragraphs.append(para)
    return paragraphs


def extract_segments_pptx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    return [para.text for para in _iter_pptx_paragraphs(prs)]


def rebuild_document_pptx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild a PPTX file with translated text replacing original segments."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    for i, para in enumerate(_iter_pptx_paragraphs(prs)):
        if i >= len(translations):
            break
        _replace_paragraph_text(para, translations[i])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# -- XLSX ---------------------------------------------------------------------


def extract_segments_xlsx(file_bytes: bytes) -> list[str]:
    """Extract translatable text segments from an XLSX file."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes))
    segments: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    segments.append(cell.value)
    return segments


def rebuild_document_xlsx(file_bytes: bytes, translations: list[str]) -> bytes:
    """Rebuild an XLSX file with translated text replacing original segments."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes))
    idx = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    if idx < len(translations):
                        cell.value = translations[idx]
                        idx += 1
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
